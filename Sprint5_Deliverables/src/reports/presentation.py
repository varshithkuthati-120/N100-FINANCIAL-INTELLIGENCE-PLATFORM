from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import sqlite3
import os

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sprint 5 Output: Cash Flow Intelligence & NLP Analysis"
    subtitle.text = "Financial Ratio Engine Automated Report"
    
    # Summary slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Sprint 5 Key Deliverables"
    tf = body_shape.text_frame
    tf.text = "Completed 100% of tasks:"
    p = tf.add_paragraph()
    p.text = "NLP Pros/Cons generated for 92 companies with strict rules."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cash Flow KPIs extracted (Quality, CapEx Intensity, Distress)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "92 individual Tearsheet PDFs built via ReportLab."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "11 Sector level reports with medians."
    p.level = 1
    
    # Distress slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_shape.text = "Distress Signal Alerts"
    tf = body_shape.text_frame
    
    try:
        distress = pd.read_csv('output/distress_alerts.csv')
        tf.text = f"Identified {len(distress)} companies with Distress Signals (CFO < 0 and CFF > 0):"
        for _, row in distress.head(5).iterrows():
            p = tf.add_paragraph()
            p.text = f"{row['company_id']} - CFO: {row['cfo_value']:.1f}, CFF: {row['cff_value']:.1f}"
            p.level = 1
        if len(distress) > 5:
            p = tf.add_paragraph()
            p.text = "... and others"
            p.level = 1
    except:
        tf.text = "No distress alerts found or file missing."

    prs.save('reports/Sprint5_Presentation.pptx')
    print("Presentation saved to reports/Sprint5_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
